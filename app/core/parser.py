"""多格式文档解析器

支持 PDF / Word / Excel / PPT / txt / md / csv / 图片(OCR via VLM)
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import tempfile
import uuid
from dataclasses import asdict
from pathlib import Path
from typing import Any

from KBzhy.app.core.document_models import DocumentElement, ParsedDocument
from KBzhy.config import PARSED_ARTIFACT_DIR

logger = logging.getLogger(__name__)


class ParsedArtifactError(ValueError):
    """Raised when a parsed artifact does not match the persisted schema."""


class DocumentParseError(ValueError):
    """Raised when a source cannot produce indexable document elements."""


class Document:
    """解析后的文档"""

    def __init__(
        self,
        content: str,
        metadata: dict[str, Any] | None = None,
    ):
        self.content = content
        self.metadata = metadata or {}
        self.doc_id: str = str(uuid.uuid4())

    def __repr__(self) -> str:
        src = self.metadata.get("source", "unknown")
        return f"Document(id={self.doc_id[:8]}, source={src}, len={len(self.content)})"


class DocumentParser:
    """多格式文档解析器"""

    SUPPORTED_TYPES: dict[str, list[str]] = {
        "pdf": [".pdf"],
        "word": [".docx"],
        "excel": [".xlsx", ".xls"],
        "ppt": [".pptx", ".ppt"],
        "text": [".txt", ".md"],
        "csv": [".csv"],
        "image": [".jpg", ".jpeg", ".png", ".bmp", ".tiff"],
    }
    _ARTIFACT_NAME_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")

    def __init__(
        self,
        vlm_model: str | None = None,
        artifact_dir: str | Path | None = None,
    ):
        self._vlm_model = vlm_model
        self._ext_map = self._build_ext_map()
        self._artifact_dir = Path(artifact_dir or PARSED_ARTIFACT_DIR).resolve()

    def _build_ext_map(self) -> dict[str, str]:
        mapping: dict[str, str] = {}
        for category, exts in self.SUPPORTED_TYPES.items():
            for ext in exts:
                mapping[ext] = category
        return mapping

    def parse(self, file_path: str | Path) -> list[Document]:
        """解析文件，返回文档列表"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = file_path.suffix.lower()
        category = self._ext_map.get(ext)
        if category is None:
            raise ValueError(f"不支持的文件格式: {ext}")

        logger.info("解析文档: %s (类型: %s)", file_path.name, category)

        parser = getattr(self, f"_parse_{category}", None)
        if parser is None:
            raise NotImplementedError(f"解析器未实现: {category}")

        docs = parser(str(file_path))
        for doc in docs:
            doc.metadata.setdefault("source", file_path.name)
            doc.metadata.setdefault("file_type", category)
            doc.metadata.setdefault("file_path", str(file_path))
        return docs

    def parse_bytes(self, content: bytes, filename: str) -> list[Document]:
        """解析字节流（上传文件用）"""
        ext = Path(filename).suffix.lower()
        category = self._ext_map.get(ext)
        if category is None:
            raise ValueError(f"不支持的文件格式: {ext}")

        parser = getattr(self, f"_parse_{category}")
        docs = parser(io.BytesIO(content)) if category in ("pdf", "image") else parser(content)
        for doc in docs:
            doc.metadata.setdefault("source", filename)
            doc.metadata.setdefault("file_type", category)
        return docs

    def parse_structured(
        self,
        file_path: str | Path,
        *,
        document_id: str,
        version: int,
        kb_id: str = "default",
    ) -> ParsedDocument:
        self._validate_identifier(document_id, "document_id")
        self._validate_identifier(kb_id, "kb_id")
        self._validate_version(version)

        source = Path(file_path)
        ext = source.suffix.lower()
        category = self._ext_map.get(ext)
        if category is None:
            raise ValueError(f"不支持的文件格式: {ext}")

        from KBzhy.app.core.parsers.pdf_parser import parse_pdf
        from KBzhy.app.core.parsers.text_parser import parse_markdown, parse_text
        from KBzhy.app.core.parsers.word_parser import parse_word

        structured_parsers = {
            ".pdf": parse_pdf,
            ".docx": parse_word,
            ".md": parse_markdown,
            ".txt": parse_text,
        }
        structured_parser = structured_parsers.get(ext)
        if structured_parser is None or not source.exists():
            documents = self.parse(source)
            elements = self._legacy_to_elements(documents, document_id)
            metadata = dict(documents[0].metadata) if documents else {}
        else:
            elements = structured_parser(source, document_id=document_id)
            metadata = {
                "source": source.name,
                "file_type": category,
                "file_path": str(source),
            }
        metadata["kb_id"] = kb_id
        return ParsedDocument(
            document_id=document_id,
            version=version,
            title=source.stem,
            language="und",
            elements=elements,
            metadata=metadata,
        )

    def save_artifact(
        self, parsed: ParsedDocument, *, artifact_name: str | None = None
    ) -> Path:
        kb_id = parsed.metadata.get("kb_id", "default")
        self._validate_identifier(kb_id, "kb_id")
        self._validate_identifier(parsed.document_id, "document_id")
        self._validate_version(parsed.version)

        target_dir = self._artifact_dir / kb_id / parsed.document_id
        self._resolve_artifact_path(target_dir)
        target_dir.mkdir(parents=True, exist_ok=True)
        target_dir = self._resolve_artifact_path(target_dir)
        stem = artifact_name if artifact_name is not None else f"v{parsed.version}"
        if not self._ARTIFACT_NAME_RE.fullmatch(stem):
            raise ValueError("artifact name is invalid")
        target = self._resolve_artifact_path(target_dir / f"{stem}.json")
        temp_path: Path | None = None
        try:
            with tempfile.NamedTemporaryFile(
                mode="w",
                encoding="utf-8",
                dir=target_dir,
                prefix=f".{target.stem}.",
                suffix=".tmp",
                delete=False,
            ) as temp_file:
                temp_path = Path(temp_file.name)
                json.dump(asdict(parsed), temp_file, ensure_ascii=False, separators=(",", ":"))
                temp_file.flush()
                os.fsync(temp_file.fileno())
            temp_path = self._resolve_artifact_path(temp_path)
            target = self._resolve_artifact_path(target)
            os.replace(temp_path, target)
        except Exception:
            if temp_path is not None:
                temp_path.unlink(missing_ok=True)
            raise
        return target

    def load_artifact(self, path: str | Path) -> ParsedDocument:
        artifact_path = self._resolve_artifact_path(path)
        try:
            data = json.loads(artifact_path.read_text(encoding="utf-8"))
            return self._decode_artifact(data)
        except (json.JSONDecodeError, KeyError, TypeError, ValueError) as exc:
            if isinstance(exc, ParsedArtifactError):
                raise
            raise ParsedArtifactError("invalid parsed artifact") from exc

    def _decode_artifact(self, data: object) -> ParsedDocument:
        if not isinstance(data, dict):
            raise ParsedArtifactError("invalid parsed artifact: root must be an object")
        for field in ("document_id", "version", "title", "language", "elements", "metadata"):
            if field not in data:
                raise ParsedArtifactError(f"invalid parsed artifact: missing {field}")

        self._validate_identifier(data["document_id"], "document_id")
        self._validate_version(data["version"])
        if not isinstance(data["elements"], list):
            raise ParsedArtifactError("invalid parsed artifact: elements must be a list")
        if not isinstance(data["metadata"], dict):
            raise ParsedArtifactError("invalid parsed artifact: metadata must be an object")
        if not isinstance(data["title"], str):
            raise ParsedArtifactError("invalid parsed artifact: title must be a string")
        if not isinstance(data["language"], str):
            raise ParsedArtifactError("invalid parsed artifact: language must be a string")

        elements = tuple(self._decode_element(element) for element in data["elements"])
        return ParsedDocument(
            document_id=data["document_id"],
            version=data["version"],
            title=data["title"],
            language=data["language"],
            elements=elements,
            metadata=data["metadata"],
        )

    @staticmethod
    def _decode_element(data: object) -> DocumentElement:
        if not isinstance(data, dict):
            raise ParsedArtifactError("invalid parsed artifact: element must be an object")
        required = ("element_id", "element_type", "text", "order")
        if any(field not in data for field in required):
            raise ParsedArtifactError("invalid parsed artifact: element fields are incomplete")
        if not isinstance(data["element_id"], str) or not data["element_id"]:
            raise ParsedArtifactError("invalid parsed artifact: element_id must be a string")
        if data["element_type"] not in {"heading", "paragraph", "list", "table", "code"}:
            raise ParsedArtifactError("invalid parsed artifact: invalid element_type")
        if not isinstance(data["text"], str):
            raise ParsedArtifactError("invalid parsed artifact: text must be a string")
        if isinstance(data["order"], bool) or not isinstance(data["order"], int):
            raise ParsedArtifactError("invalid parsed artifact: order must be an integer")
        page = data.get("page")
        if page is not None and (isinstance(page, bool) or not isinstance(page, int)):
            raise ParsedArtifactError("invalid parsed artifact: page must be an integer")
        section_path = data.get("section_path", [])
        if not isinstance(section_path, list) or not all(isinstance(part, str) for part in section_path):
            raise ParsedArtifactError("invalid parsed artifact: section_path must be a string list")
        bounding_box = data.get("bounding_box")
        if bounding_box is not None and (
            not isinstance(bounding_box, dict)
            or not all(
                isinstance(key, str)
                and not isinstance(value, bool)
                and isinstance(value, (int, float))
                for key, value in bounding_box.items()
            )
        ):
            raise ParsedArtifactError("invalid parsed artifact: bounding_box must contain numbers")
        metadata = data.get("metadata", {})
        if not isinstance(metadata, dict):
            raise ParsedArtifactError("invalid parsed artifact: element metadata must be an object")
        return DocumentElement(
            element_id=data["element_id"],
            element_type=data["element_type"],
            text=data["text"],
            order=data["order"],
            page=page,
            section_path=tuple(section_path),
            bounding_box=bounding_box,
            metadata=metadata,
        )

    def _resolve_artifact_path(self, path: str | Path) -> Path:
        resolved = Path(path).resolve()
        try:
            resolved.relative_to(self._artifact_dir)
        except ValueError as exc:
            raise ValueError("artifact path must stay within artifact_dir") from exc
        return resolved

    @staticmethod
    def _legacy_to_elements(
        documents: list[Document],
        document_id: str,
    ) -> tuple[DocumentElement, ...]:
        elements: list[DocumentElement] = []
        for order, document in enumerate(documents):
            metadata = dict(document.metadata)
            is_table = metadata.get("file_type") in {"excel", "csv"} or "sheet" in metadata
            elements.append(
                DocumentElement(
                    element_id=f"{document_id}:legacy:{order}",
                    element_type="table" if is_table else "paragraph",
                    text=document.content,
                    order=order,
                    page=metadata.get("page"),
                    section_path=DocumentParser._normalize_section_path(
                        metadata.get("section_path")
                    ),
                    bounding_box=metadata.get("bounding_box"),
                    metadata=metadata,
                )
            )
        return tuple(elements)

    @staticmethod
    def _normalize_section_path(value: object) -> tuple[str, ...]:
        if value is None:
            return ()
        if isinstance(value, str):
            return (value,)
        if isinstance(value, (list, tuple)) and all(isinstance(part, str) for part in value):
            return tuple(value)
        raise ParsedArtifactError("invalid legacy section_path")

    @staticmethod
    def _validate_identifier(value: object, field: str) -> None:
        if not isinstance(value, str) or not value or value in {".", ".."}:
            raise ValueError(f"{field} must be a safe path segment")
        if re.fullmatch(r"[\w.-]+", value) is None:
            raise ValueError(f"{field} must be a safe path segment")

    @staticmethod
    def _validate_version(version: object) -> None:
        if isinstance(version, bool) or not isinstance(version, int) or version < 1:
            raise ValueError("version must be a positive integer")

    # ── PDF ────────────────────────────────────

    def _parse_pdf(self, source: str | io.BytesIO) -> list[Document]:
        import fitz

        docs: list[Document] = []
        if isinstance(source, str):
            pdf = fitz.open(source)
        else:
            pdf = fitz.open(stream=source.read(), filetype="pdf")

        for page_num in range(len(pdf)):
            page = pdf[page_num]
            text = page.get_text().strip()
            if text:
                docs.append(Document(
                    content=text,
                    metadata={"page": page_num + 1},
                ))
        pdf.close()
        return docs if docs else [Document(content="[PDF 无法提取文字，可能是扫描件]", metadata={"page": 0})]

    # ── Word ───────────────────────────────────

    def _parse_word(self, source: str | bytes) -> list[Document]:
        from docx import Document as DocxDocument

        if isinstance(source, str):
            docx = DocxDocument(source)
        else:
            docx = DocxDocument(io.BytesIO(source))

        full_text: list[str] = []
        for para in docx.paragraphs:
            full_text.append(para.text)

        for table in docx.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                full_text.append(" | ".join(cells))

        return [Document(content="\n".join(full_text))]

    # ── Excel ──────────────────────────────────

    def _parse_excel(self, source: str | bytes) -> list[Document]:
        import openpyxl

        if isinstance(source, str):
            wb = openpyxl.load_workbook(source, data_only=True)
        else:
            wb = openpyxl.load_workbook(io.BytesIO(source), data_only=True)

        docs: list[Document] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                docs.append(Document(
                    content="\n".join(rows),
                    metadata={"sheet": sheet_name},
                ))
        wb.close()
        return docs if docs else [Document(content="[Excel 文件为空]")]

    # ── PPT ────────────────────────────────────

    def _parse_ppt(self, source: str | bytes) -> list[Document]:
        from pptx import Presentation

        if isinstance(source, str):
            prs = Presentation(source)
        else:
            prs = Presentation(io.BytesIO(source))

        docs: list[Document] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                docs.append(Document(
                    content="\n".join(texts),
                    metadata={"slide": slide_num},
                ))
        return docs if docs else [Document(content="[PPT 无文字内容]")]

    # ── 纯文本 ─────────────────────────────────

    @staticmethod
    def _read_with_fallback(file_path: str) -> str:
        """读取文本文件，UTF-8 → GBK → GB18030 → latin-1 依次尝试"""
        encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312", "gb18030", "latin-1"]
        for enc in encodings:
            try:
                return Path(file_path).read_text(encoding=enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return Path(file_path).read_text(encoding="utf-8", errors="replace")

    @staticmethod
    def _decode_with_fallback(data: bytes) -> str:
        """解码字节，UTF-8 → GBK → GB18030 → latin-1 依次尝试"""
        encodings = ["utf-8", "utf-8-sig", "gbk", "gb2312", "gb18030", "latin-1"]
        for enc in encodings:
            try:
                return data.decode(enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return data.decode("utf-8", errors="replace")

    def _parse_text(self, source: str | bytes) -> list[Document]:
        if isinstance(source, str):
            content = self._read_with_fallback(source)
        else:
            content = self._decode_with_fallback(source)
        return [Document(content=content)]

    # ── CSV ────────────────────────────────────

    def _parse_csv(self, source: str | bytes) -> list[Document]:
        import csv

        if isinstance(source, str):
            text = self._read_with_fallback(source)
        else:
            text = self._decode_with_fallback(source)

        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(row) for row in reader if any(c.strip() for c in row)]
        return [Document(content="\n".join(rows))]

    # ── 图片 (OCR via VLM) ─────────────────────

    def _parse_image(self, source: str | io.BytesIO) -> list[Document]:
        import base64

        if isinstance(source, str):
            with open(source, "rb") as f:
                img_data = base64.b64encode(f.read()).decode()
            source_name = os.path.basename(source)
        else:
            img_data = base64.b64encode(source.read()).decode()
            source_name = "uploaded_image"

        from KBzhy.config import API_KEY, API_BASE, LLM_MODEL

        if not API_KEY or API_KEY == "your-api-key":
            return [Document(content="[图片 OCR 需要配置 API Key]")]

        model = self._vlm_model or LLM_MODEL

        try:
            import httpx
            resp = httpx.post(
                f"{API_BASE.rstrip('/')}/chat/completions",
                headers={
                    "Authorization": f"Bearer {API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:image/jpeg;base64,{img_data}"},
                                },
                                {
                                    "type": "text",
                                    "text": "请提取并输出这张图片中的所有文字内容，保持原文格式（表格、列表等结构）。只输出文字，不要添加额外说明。",
                                },
                            ],
                        }
                    ],
                    "max_tokens": 4096,
                },
                timeout=60.0,
            )
            resp.raise_for_status()
            data = resp.json()
            text = data["choices"][0]["message"]["content"]
            return [Document(content=text, metadata={"ocr_method": "vlm", "source": source_name})]
        except Exception as exc:
            logger.error("VLM OCR 失败: %s", exc)
            return [Document(content=f"[图片 OCR 失败: {exc}]", metadata={"source": source_name})]
